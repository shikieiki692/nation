#!/usr/bin/env python3
"""
原子结构（三）PPT 生成脚本
基于 python-pptx，覆盖第八讲全部 6 个 §，
参考超级充实版内容 + 教学洞察 + 原讲义填空答案。
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

# ============================================================
# 0. CONSTANTS
# ============================================================

# Colors
DARK_BLUE = RGBColor(0x1E, 0x40, 0xAF)
BRIGHT_BLUE = RGBColor(0x3B, 0x82, 0xF6)
MEDIUM_BLUE = RGBColor(0x60, 0xA5, 0xFA)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
DARK_GRAY = RGBColor(0x37, 0x41, 0x51)
GRAY = RGBColor(0x6B, 0x72, 0x80)
LIGHT_GRAY = RGBColor(0xF3, 0xF4, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xDC, 0x26, 0x26)
GREEN = RGBColor(0x05, 0x92, 0x2B)
LIGHT_AMBER = RGBColor(0xFE, 0xF3, 0xC7)
LIGHT_RED = RGBColor(0xFE, 0xE2, 0xE2)
LIGHT_BLUE = RGBColor(0xDB, 0xEA, 0xFE)
LIGHT_GREEN = RGBColor(0xD1, 0xFA, 0xE5)
DARK_BLUE_BG = RGBColor(0x1E, 0x3A, 0x5F)  # for section dividers

# Fonts
CN_TITLE = 'SimHei'
CN_BODY = 'SimSun'
EN_FONT = 'Times New Roman'

# Slide dimensions (16:9)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Paths
VAULT_ROOT = r'C:\Obsidion\妙妙屋'
IMG_DIR = os.path.join(VAULT_ROOT, r'原子结构讲义\第八讲原子结构（三）_images')
OUTPUT_DIR = os.path.join(VAULT_ROOT, r'04-课件\PPT')
OUTPUT_FILE = os.path.join(OUTPUT_DIR, '原子结构-第八讲-授课版.pptx')

# Margins
ML = 0.7   # margin left
MR = 0.7   # margin right
MT = 0.5   # margin top
CONTENT_W = 13.333 - ML - MR  # ~11.93

# ============================================================
# 1. HELPER FUNCTIONS
# ============================================================

def set_font(run, cn_font=CN_BODY, size=Pt(14), bold=False, color=DARK_GRAY, italic=False):
    """Set both western and east-asian fonts on a run."""
    run.font.name = EN_FONT
    run.font.size = size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    rPr.set(qn('w:eastAsia'), cn_font)


def add_paragraph(tf, text, cn_font=CN_BODY, size=Pt(14), bold=False, color=DARK_GRAY,
                  alignment=PP_ALIGN.LEFT, space_before=0, space_after=0, italic=False):
    """Add a paragraph to a text frame."""
    p = tf.add_paragraph()
    p.alignment = alignment
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    run = p.add_run()
    run.text = text
    set_font(run, cn_font, size, bold, color, italic)
    return p


def add_textbox(slide, left, top, width, height, text="", cn_font=CN_BODY,
                size=Pt(14), bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT,
                word_wrap=True, italic=False):
    """Add a text box with one paragraph."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    set_font(run, cn_font, size, bold, color, italic)
    return txBox


def add_card(slide, left, top, width, height, bg_color=LIGHT_GRAY):
    """Add a rounded rectangle card."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    # Reduce corner rounding
    shape.adjustments[0] = 0.05
    return shape


def add_callout_box(slide, left, top, width, height, icon, title, body_text, bg_color):
    """Add a callout card with icon header and body."""
    card = add_card(slide, left, top, width, height, bg_color)
    tf = card.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    # Icon + Title line
    run = tf.paragraphs[0].add_run()
    run.text = f"{icon} {title}"
    set_font(run, CN_TITLE, Pt(13), True, DARK_GRAY)
    # Body
    p = add_paragraph(tf, body_text, CN_BODY, Pt(12), color=DARK_GRAY, space_before=4)
    return card


def add_section_divider(slide, section_num, title, subtitle=""):
    """Add a full-slide section divider with dark blue background."""
    # Background
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BLUE_BG
    bg.line.fill.background()
    # Decorative bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.8), Inches(3.2), Inches(1.2), Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = AMBER
    bar.line.fill.background()
    # Section number
    add_textbox(slide, 0.8, 2.2, 5, 0.8, f"§{section_num}",
                CN_TITLE, Pt(20), False, AMBER, PP_ALIGN.LEFT)
    # Title
    add_textbox(slide, 0.8, 3.4, 10, 1.2, title,
                CN_TITLE, Pt(32), True, WHITE, PP_ALIGN.LEFT)
    if subtitle:
        add_textbox(slide, 0.8, 4.6, 10, 0.8, subtitle,
                    CN_BODY, Pt(16), False, RGBColor(0xBF, 0xDB, 0xFE), PP_ALIGN.LEFT)


def add_page_number(slide, num, total):
    """Add page number in bottom right."""
    add_textbox(slide, 12.0, 7.0, 1.2, 0.4, f"{num}/{total}",
                CN_BODY, Pt(9), False, GRAY, PP_ALIGN.RIGHT)


def add_slide_title(slide, title, subtitle=""):
    """Add standard slide title bar at top."""
    # Title background bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), SLIDE_W, Inches(0.9))
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_BLUE
    bar.line.fill.background()
    # Title text
    add_textbox(slide, ML, 0.12, CONTENT_W, 0.6, title,
                CN_TITLE, Pt(22), True, WHITE, PP_ALIGN.LEFT)
    if subtitle:
        add_textbox(slide, ML, 0.55, CONTENT_W, 0.35, subtitle,
                    CN_BODY, Pt(11), False, RGBColor(0xBF, 0xDB, 0xFE), PP_ALIGN.LEFT)
    # Bottom accent line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0.9), SLIDE_W, Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = AMBER
    line.line.fill.background()


def make_content_slide(prs, title, body_elements, page_num=None, total=None):
    """Make a standard content slide with title + body text.
    body_elements: list of (text, cn_font, size, bold, color, alignment) tuples
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_slide_title(slide, title)
    # Content area
    y = 1.2
    for elem in body_elements:
        text = elem[0]
        font = elem[1] if len(elem) > 1 else CN_BODY
        size = elem[2] if len(elem) > 2 else Pt(14)
        bold = elem[3] if len(elem) > 3 else False
        color = elem[4] if len(elem) > 4 else DARK_GRAY
        align = elem[5] if len(elem) > 5 else PP_ALIGN.LEFT
        tb = add_textbox(slide, ML, y, CONTENT_W, 0.6, text, font, size, bold, color, align)
        # Estimate height based on text length
        lines = max(1, len(text) / 60)
        y += 0.08 + lines * (size.pt / 18)  # approximate
    return slide


# ============================================================
# 2. BUILD PRESENTATION
# ============================================================

def build_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Track page numbers
    total = 0
    slides_meta = []  # not used directly, counter below

    # ========================
    # SLIDE 1: COVER
    # ========================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Dark blue background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid(); bg.fill.fore_color.rgb = DARK_BLUE_BG; bg.line.fill.background()
    # Decorative elements
    for i, (w, h, x, y, c) in enumerate([
        (3.0, 0.06, 0.8, 2.8, AMBER),
        (0.06, 2.0, 1.6, 3.5, BRIGHT_BLUE),
    ]):
        d = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        d.fill.solid(); d.fill.fore_color.rgb = c; d.line.fill.background()
    # Main title
    add_textbox(slide, 2.0, 2.6, 10, 1.2, "原子结构（三）",
                CN_TITLE, Pt(42), True, WHITE, PP_ALIGN.LEFT)
    # Subtitle
    add_textbox(slide, 2.0, 3.8, 10, 0.8, "从量子力学描述到核外电子排布",
                CN_BODY, Pt(20), False, RGBColor(0xBF, 0xDB, 0xFE), PP_ALIGN.LEFT)
    # Badges
    badges = "📖 第八讲  ·  🧪 初二竞赛班  ·  ⏱ 第9课收官"
    add_textbox(slide, 2.0, 5.0, 8, 0.5, badges,
                CN_BODY, Pt(13), False, AMBER, PP_ALIGN.LEFT)
    # Footer info
    add_textbox(slide, 2.0, 6.2, 8, 0.4, "参考：第八讲原子结构（三）· 原子结构超级充实版 · 教学洞察-原子结构",
                CN_BODY, Pt(9), False, GRAY, PP_ALIGN.LEFT)
    total += 1

    # ========================
    # SLIDE 2: TABLE OF CONTENTS
    # ========================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "目  录", "第八讲原子结构（三）· 共6章")
    sections = [
        ("§一", "原子结构发现史", "从 Dalton 实心球到薛定谔波动方程"),
        ("§二", "微观粒子运动的基本特征", "波粒二象性 · 德布罗意 · 测不准原理"),
        ("§三", "原子结构的量子力学描述", "薛定谔方程 · 四量子数 · 概率密度与电子云"),
        ("§四", "原子核外电子排布", "构造原理 · 屏蔽与钻穿 · Slater规则 · 三原则"),
        ("§五", "核外电子排布（续）", "离子排布 · 周期律趋势 · 电离能反常"),
        ("§六", "原子核衰变", "α/β/γ衰变 · 核反应方程式 · 半衰期"),
    ]
    y = 1.3
    for i, (num, title, desc) in enumerate(sections):
        # Color coding
        colors = [DARK_BLUE, BRIGHT_BLUE, MEDIUM_BLUE, AMBER, GREEN, RGBColor(0x7C, 0x3A, 0xED)]
        c = colors[i]
        # Number box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(ML + 0.2), Inches(y), Inches(0.8), Inches(0.7))
        box.fill.solid(); box.fill.fore_color.rgb = c; box.line.fill.background()
        box_tf = box.text_frame; box_tf.word_wrap = True
        box_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        r = box_tf.paragraphs[0].add_run(); r.text = num
        set_font(r, CN_TITLE, Pt(16), True, WHITE)
        # Title
        add_textbox(slide, ML + 1.3, y - 0.02, 4, 0.45, title,
                    CN_TITLE, Pt(16), True, DARK_GRAY, PP_ALIGN.LEFT)
        add_textbox(slide, ML + 1.3, y + 0.35, 8, 0.35, desc,
                    CN_BODY, Pt(11), False, GRAY, PP_ALIGN.LEFT)
        y += 0.95
    total += 1

    # ================================================================
    # SECTION 1: 原子结构发现史 (Slides 3-8)
    # ================================================================

    # --- Divider ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_divider(slide, "一", "原子结构发现史",
                        "Dalton → Thomson → Rutherford → Bohr → Schrödinger")
    total += 1

    # --- Slide: 五大模型时间线 ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "原子模型演变的五大里程碑", "每一次突破都来自实验异常 → 旧模型失效 → 新假设提出")
    models = [
        ("1803 Dalton", "实心球模型", "原子不可再分", "倍比定律"),
        ("1903 Thomson", "葡萄干布丁", "正电荷均匀分布", "发现电子"),
        ("1911 Rutherford", "核式模型", "核（正+质量）+ 核外电子", "α粒子散射"),
        ("1913 Bohr", "行星模型", "定态 + 角动量量子化", "氢原子线状光谱"),
        ("1926 Schrödinger", "量子力学模型", "ψ波函数描述，|ψ|²概率密度", "德布罗意 + 电子衍射"),
    ]
    y = 1.2
    cols = [0.7, 1.8, 4.5, 11.5]  # x positions for columns
    # Header row
    headers = ["年份", "模型", "核心假设", "实验依据"]
    for j, (hdr, x) in enumerate(zip(headers, cols)):
        tb = add_textbox(slide, x, y, 2.5, 0.35, hdr, CN_TITLE, Pt(11), True, DARK_BLUE)
    y += 0.4
    for i, (year, name, hypo, exp) in enumerate(models):
        bg_color = LIGHT_GRAY if i % 2 == 0 else WHITE
        row_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                         Inches(ML-0.1), Inches(y), Inches(CONTENT_W+0.2), Inches(0.55))
        row_bg.fill.solid(); row_bg.fill.fore_color.rgb = bg_color; row_bg.line.fill.background()
        add_textbox(slide, cols[0], y+0.05, cols[1]-cols[0]-0.1, 0.4, year,
                    CN_BODY, Pt(11), True, DARK_BLUE)
        add_textbox(slide, cols[1], y+0.05, cols[2]-cols[1]-0.1, 0.4, name,
                    CN_TITLE, Pt(12), True, DARK_GRAY)
        add_textbox(slide, cols[2], y+0.05, cols[3]-cols[2]-0.3, 0.4, hypo,
                    CN_BODY, Pt(11), False, DARK_GRAY)
        add_textbox(slide, cols[3], y+0.05, 1.8, 0.4, exp,
                    CN_BODY, Pt(10), False, GRAY)
        y += 0.55
    total += 1

    # --- Slide: α散射实验与Rutherford核式模型 ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "卢瑟福α粒子散射实验（1911）", "1/8000 的粒子偏转 > 90°，汤姆孙模型无法解释")
    # Key facts left
    y = 1.3
    add_textbox(slide, ML, y, 5.5, 0.4, "🔬 实验装置与方法", CN_TITLE, Pt(14), True, DARK_BLUE)
    y += 0.5
    points = [
        "用准直的α射线轰击厚度为微米的金箔",
        "绝大多数α粒子照直穿过（偏转很小）",
        "少数发生较大偏转，约 1/8000 偏转 > 90°",
        "观察到偏转角等于 150° 的散射",
        "结论：原子大部分是空的，正电荷集中在极小核区",
    ]
    for pt in points:
        add_textbox(slide, ML + 0.3, y, 5.2, 0.35, f"• {pt}", CN_BODY, Pt(12), False, DARK_GRAY)
        y += 0.35
    # Right card: Rutherford model
    y = 1.3
    card = add_card(slide, 6.8, y, 5.8, 2.8, LIGHT_BLUE)
    tf = card.text_frame; tf.word_wrap = True
    run = tf.paragraphs[0].add_run(); run.text = "⭐ 卢瑟福核式模型要点"
    set_font(run, CN_TITLE, Pt(13), True, DARK_BLUE)
    items = [
        "每个原子中心有一个带正电荷的原子核",
        "核外电子绕核高速运动",
        "原子核体积很小但集中了几乎全部质量",
        "核正电荷 = 核外负电荷 → 原子呈电中性",
    ]
    for item in items:
        p = add_paragraph(tf, f"  ✓ {item}", CN_BODY, Pt(11), color=DARK_GRAY, space_before=3)
    total += 1

    # --- Slide: Bohr模型三项假设 ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "玻尔模型的三项假设（1913）", "旧量子论的巅峰——用 Planck 量子论解释氢原子光谱")
    y = 1.3
    # Three hypothesis cards
    hypos = [
        ("① 定态假设", "电子只能在特定半径的轨道上运动\n且不辐射能量 → 称\"定态\"",
         DARK_BLUE),
        ("② 角动量量子化", "L = n·h/2π, n=1,2,3,...\n只有角动量为 h/2π 整数倍的轨道才允许",
         BRIGHT_BLUE),
        ("③ 跃迁假设", "电子从高能态 E₂ 跃迁到低能态 E₁ 时\n放出光子：ΔE = E₂ - E₁ = hν",
         AMBER),
    ]
    card_w = 3.5
    gap = 0.4
    start_x = ML
    for i, (title, body, clr) in enumerate(hypos):
        x = start_x + i * (card_w + gap)
        card = add_card(slide, x, y, card_w, 2.0, LIGHT_GRAY)
        tf = card.text_frame; tf.word_wrap = True
        r = tf.paragraphs[0].add_run(); r.text = title
        set_font(r, CN_TITLE, Pt(14), True, clr)
        p = add_paragraph(tf, body, CN_BODY, Pt(12), color=DARK_GRAY, space_before=6)
    # Bottom formula
    y = 3.6
    add_textbox(slide, ML, y, 5, 0.5, "轨道半径：rₙ = a₀·n² （a₀ = 52.9 pm）",
                CN_BODY, Pt(13), False, DARK_GRAY)
    add_textbox(slide, ML + 5.5, y, 5, 0.5, "电子能量：Eₙ = -13.6/n² eV",
                CN_BODY, Pt(13), True, DARK_BLUE)
    # 教学洞察 callout
    add_callout_box(slide, ML, 4.3, CONTENT_W, 1.0, "🧠", "教学洞察：经典物理的第一个耳光",
                    "\"如果电子绕核做圆周运动，根据经典电磁理论——加速运动的电荷会辐射电磁波，能量损失导致电子坠入原子核。但原子明明是稳定的！\" — 质心L2",
                    LIGHT_AMBER)
    total += 1

    # --- Slide: 氢原子光谱与Balmer系 ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "氢原子光谱——量子化的实验证据", "线状光谱 → 能量量子化 → 能级分立")
    # Left: spectrum table
    y = 1.3
    add_textbox(slide, ML, y, 4, 0.4, "氢原子五组线系", CN_TITLE, Pt(14), True, DARK_BLUE)
    y += 0.45
    lines_data = [
        ("Lyman系", "n=1", "紫外", "1/1² - 1/n²"),
        ("Balmer系", "n=2", "可见光", "1/2² - 1/n²"),
        ("Paschen系", "n=3", "红外", "1/3² - 1/n²"),
        ("Brackett系", "n=4", "红外", "1/4² - 1/n²"),
        ("Pfund系", "n=5", "远红外", "1/5² - 1/n²"),
    ]
    for name, low, region, formula in lines_data:
        bg_c = LIGHT_GRAY if lines_data.index((name, low, region, formula)) % 2 == 0 else WHITE
        row_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            Inches(ML), Inches(y), Inches(5.5), Inches(0.35))
        row_bg.fill.solid(); row_bg.fill.fore_color.rgb = bg_c; row_bg.line.fill.background()
        add_textbox(slide, ML+0.1, y+0.02, 1.2, 0.3, name, CN_TITLE, Pt(11), True, DARK_BLUE)
        add_textbox(slide, ML+1.3, y+0.02, 1.0, 0.3, low, CN_BODY, Pt(11), False, DARK_GRAY)
        add_textbox(slide, ML+2.3, y+0.02, 1.0, 0.3, region, CN_BODY, Pt(10), False, GRAY)
        add_textbox(slide, ML+3.3, y+0.02, 2.2, 0.3, f"ν̃ = R_H({formula})",
                    CN_BODY, Pt(9), False, GRAY)
        y += 0.35
    # Right: key formula
    y = 1.3
    card = add_card(slide, 7.0, y, 5.5, 3.2, LIGHT_BLUE)
    tf = card.text_frame; tf.word_wrap = True
    r = tf.paragraphs[0].add_run(); r.text = "📐 核心公式"
    set_font(r, CN_TITLE, Pt(13), True, DARK_BLUE)
    formulas = [
        "",
        "Rydberg 公式：",
        "  ν̃ = R_H (1/n₁² - 1/n₂²)",
        "",
        "Rydberg 常数：",
        "  R_H = 1.097 × 10⁵ cm⁻¹",
        "",
        "Balmer 系 Hα 线（n₁=2, n₂=3）：",
        "  λ = 656.3 nm（可见光红色）",
    ]
    for f in formulas:
        bold = any(k in f for k in ["Rydberg", "Balmer"])
        add_paragraph(tf, f, CN_BODY, Pt(12), bold, DARK_GRAY, space_before=2)
    # Bottom 🧠
    add_callout_box(slide, ML, 5.0, CONTENT_W, 1.0, "🧠", "教学洞察：从连续到离散的认知跳跃",
                    "经典电磁学预言原子光谱是连续谱——像白炽灯。但实际是线状的！只有特定波长出现。这说明能级是量子化的。光谱是\"通往量子化的钥匙\"。",
                    LIGHT_AMBER)
    # Embed: Hα计算
    y_calc = 5.0
    total += 1

    # --- Slide: 嵌入式计算：Hα线波长 ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "📝 嵌入式计算：Balmer系 Hα 线波长", "电子从 n=3 → n=2 跃迁")
    y = 1.5
    add_textbox(slide, ML, y, 10, 0.5, "已知 R_H = 1.097 × 10⁵ cm⁻¹，计算 Hα 线的波长（nm）",
                CN_BODY, Pt(14), True, DARK_GRAY)
    y += 0.7
    add_textbox(slide, ML, y, 10, 0.4, "→ ν̃ = R_H (1/2² - 1/3²) = 1.097×10⁵ × (1/4 - 1/9) = 1.097×10⁵ × 5/36",
                CN_BODY, Pt(13), False, DARK_GRAY)
    y += 0.45
    add_textbox(slide, ML, y, 10, 0.4, "→ ν̃ = 15236 cm⁻¹",
                CN_BODY, Pt(13), True, DARK_BLUE)
    y += 0.45
    add_textbox(slide, ML, y, 10, 0.4, "→ λ = 1/ν̃ = 6.563×10⁻⁵ cm = 656.3 nm ✓",
                CN_BODY, Pt(13), True, DARK_BLUE)
    y += 0.6
    add_textbox(slide, ML, y, 10, 0.4, "验证：656.3 nm 对应可见光红色，与实验观测完全一致！",
                CN_BODY, Pt(13), False, GREEN)
    # Bottom: 局限
    y = 3.8
    card = add_card(slide, ML, y, CONTENT_W, 2.5, LIGHT_RED)
    tf = card.text_frame; tf.word_wrap = True
    r = tf.paragraphs[0].add_run(); r.text = "⚠️ 玻尔模型的局限"
    set_font(r, CN_TITLE, Pt(13), True, RED)
    limits = [
        "不能解释氢原子光谱精细结构（每一条谱线实际分裂为两条）",
        "不能解释多电子原子（He 即出现显著偏差）",
        "定态假设本身与经典电磁理论相悖",
        "无法解释 Zeeman 效应（磁场中谱线分裂）",
        "→ \"波尔模型是一座桥——走过它到达量子力学，但桥本身不是目的地\"",
    ]
    for lim in limits:
        icon = "✗ " if lim.startswith("不能") or lim.startswith("定态") or lim.startswith("无法") else "→"
        bold = lim.startswith("→")
        add_paragraph(tf, f"  {icon} {lim}", CN_BODY, Pt(11), bold, DARK_GRAY, space_before=2)
    total += 1

    # ================================================================
    # SECTION 2: 微观粒子运动的基本特征 (Slides 9-13)
    # ================================================================

    # --- Divider ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_divider(slide, "二", "微观粒子运动的基本特征",
                        "波粒二象性 · 德布罗意关系式 · 海森堡测不准原理")
    total += 1

    # --- Slide: 波粒二象性 ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "波粒二象性——微观粒子的双重身份", "光：干涉衍射=波动性，光电效应=微粒性")
    y = 1.3
    # Light duality
    add_textbox(slide, ML, y, 5, 0.4, "💡 光的波粒二象性", CN_TITLE, Pt(14), True, DARK_BLUE)
    y += 0.5
    card = add_card(slide, ML, y, 5.5, 1.8, LIGHT_GRAY)
    tf = card.text_frame; tf.word_wrap = True
    r = tf.paragraphs[0].add_run(); r.text = "波动性"
    set_font(r, CN_TITLE, Pt(12), True, BRIGHT_BLUE)
    add_paragraph(tf, "  干涉、衍射；有波长 λ、频率 ν；传播方向波峰波谷",
                  CN_BODY, Pt(11), color=DARK_GRAY, space_before=3)
    p = add_paragraph(tf, "", CN_BODY, Pt(6), color=DARK_GRAY)
    r2 = tf.paragraphs[-1].add_run(); r2.text = "微粒性"
    set_font(r2, CN_TITLE, Pt(12), True, AMBER)
    add_paragraph(tf, "  发射、吸收、光电效应；E = hν，光量子能量不连续",
                  CN_BODY, Pt(11), color=DARK_GRAY, space_before=3)
    # Right: De Broglie
    y = 1.3
    add_textbox(slide, 7.0, y, 5, 0.4, "⚛ 电子的波粒二象性", CN_TITLE, Pt(14), True, DARK_BLUE)
    y += 0.5
    card2 = add_card(slide, 7.0, y, 5.5, 1.8, LIGHT_BLUE)
    tf2 = card2.text_frame; tf2.word_wrap = True
    r = tf2.paragraphs[0].add_run(); r.text = "1924 德布罗意天才猜想"
    set_font(r, CN_TITLE, Pt(12), True, DARK_BLUE)
    add_paragraph(tf2, "  λ = h/p = h/(mv)", CN_BODY, Pt(14), True, DARK_GRAY, space_before=6)
    add_paragraph(tf2, "  m 越大、v 越快 → λ 越短 → 波动性越不明显",
                  CN_BODY, Pt(11), color=DARK_GRAY, space_before=4)
    add_paragraph(tf2, "  电子质量极小(9.1×10⁻³¹ kg)，在原子尺度下波长与原子大小相当",
                  CN_BODY, Pt(11), color=DARK_GRAY, space_before=4)
    # Bottom: experimental verification
    y = 3.7
    add_textbox(slide, ML, y, 10, 0.4, "🧪 实验验证：1927 Davisson & Germer 电子衍射实验",
                CN_TITLE, Pt(13), True, DARK_BLUE)
    add_textbox(slide, ML, y + 0.45, 10, 0.4,
                "电子在 Ni 晶体表面产生清晰衍射条纹——与 X 射线衍射类似，直接证明电子具有波动性",
                CN_BODY, Pt(12), False, DARK_GRAY)
    total += 1

    # --- Slide: 海森堡测不准原理 ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "海森堡测不准原理（1927）", "不能同时准确测定微观粒子的位置和动量")
    y = 1.5
    # Main formula card
    card = add_card(slide, ML, y, 6.0, 2.2, LIGHT_AMBER)
    tf = card.text_frame; tf.word_wrap = True
    r = tf.paragraphs[0].add_run(); r.text = "📐 数学表达式"
    set_font(r, CN_TITLE, Pt(14), True, DARK_BLUE)
    add_paragraph(tf, "", CN_BODY, Pt(4))
    add_paragraph(tf, "    Δx · Δp ≥ h / (4π)", CN_BODY, Pt(20), True, DARK_BLUE, PP_ALIGN.CENTER)
    add_paragraph(tf, "", CN_BODY, Pt(4))
    add_paragraph(tf, "    Δx: 位置不确定度    Δp: 动量不确定度    h: Planck 常数",
                  CN_BODY, Pt(11), False, GRAY, PP_ALIGN.CENTER)
    # Right: explanation
    add_textbox(slide, 7.2, y, 5.5, 1.5,
                "\"你想看电子在哪儿——用光子照它。但光子有动量——一照就把电子踢走了。测得越精确（Δx 小）→ 光子波长越短 → 动量越大 → Δp 越大。你永远不能同时精确知道位置和动量。\"",
                CN_BODY, Pt(12), False, DARK_GRAY, italic=True)
    y += 2.8
    # Significance
    card2 = add_card(slide, ML, y, CONTENT_W, 1.8, LIGHT_RED)
    tf2 = card2.text_frame; tf2.word_wrap = True
    r = tf2.paragraphs[0].add_run(); r.text = "⚠️ 重要推论：轨道概念在原子尺度失去意义"
    set_font(r, CN_TITLE, Pt(13), True, RED)
    add_paragraph(tf2, "  • 不可能同时知道电子的精确位置和精确速度",
                  CN_BODY, Pt(12), color=DARK_GRAY, space_before=4)
    add_paragraph(tf2, "  • 不存在 Rutherford 和 Bohr 模型中行星绕太阳那样的电子轨道",
                  CN_BODY, Pt(12), color=DARK_GRAY, space_before=2)
    add_paragraph(tf2, "  • 只能确定电子在某空间附近出现的概率（统计规律）",
                  CN_BODY, Pt(12), color=DARK_GRAY, space_before=2)
    total += 1

    # --- Slide: 嵌入式计算：测不准原理算例 ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "📝 嵌入式计算：原子中不存在电子\"轨道\"", "定量证明轨道概念的失效")
    y = 1.5
    add_textbox(slide, ML, y, 10, 0.5, "电子在原子中的运动范围约 r ≈ 10⁻¹⁰ m。若位置误差 Δx = 10⁻¹⁰ m，",
                CN_BODY, Pt(14), False, DARK_GRAY)
    y += 0.45
    add_textbox(slide, ML, y, 10, 0.5, "计算电子速度的不确定度 Δv = ?",
                CN_BODY, Pt(14), True, DARK_GRAY)
    y += 0.6
    # Step-by-step solution
    steps = [
        ("Step 1", "Δp ≥ h / (4π·Δx)"),
        ("Step 2", "Δp ≥ 6.626×10⁻³⁴ / (4π × 10⁻¹⁰) ≈ 5.27×10⁻²⁵ kg·m/s"),
        ("Step 3", "Δv ≥ Δp / mₑ = 5.27×10⁻²⁵ / 9.109×10⁻³¹"),
        ("Result", "Δv ≥ 5.8×10⁵ m/s 🚀"),
    ]
    for label, expr in steps:
        bg_c = LIGHT_BLUE if label == "Result" else (LIGHT_GRAY if label.startswith("Step") and int(label[-1]) % 2 == 1 else WHITE)
        card = add_card(slide, ML, y, CONTENT_W, 0.5, bg_c)
        tf = card.text_frame; tf.word_wrap = True
        r = tf.paragraphs[0].add_run()
        prefix = f"  >> {expr}" if not label.startswith("Step") else f"  {label}: {expr}"
        r.text = prefix
        set_font(r, CN_BODY, Pt(14), label in ["Result"], DARK_BLUE if label == "Result" else DARK_GRAY)
        y += 0.55
    # Conclusion
    y += 0.2
    card = add_card(slide, ML, y, CONTENT_W, 1.0, LIGHT_AMBER)
    tf = card.text_frame; tf.word_wrap = True
    r = tf.paragraphs[0].add_run(); r.text = "🧠 结论：Δv ≈ 5.8×10⁵ m/s，与电子在原子中的运动速度（~10⁶ m/s）同一量级！"
    set_font(r, CN_TITLE, Pt(13), True, DARK_BLUE)
    add_paragraph(tf, "这意味着我们完全无法确定电子\"在轨道上某点的速度\"——\"轨道\"概念在原子尺度失去了意义。",
                  CN_BODY, Pt(12), color=DARK_GRAY, space_before=4)
    total += 1

    # --- Slide: 从轨道到概率 ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "从\"轨道\"到\"概率\"——思维升级", "放弃'轨道'，接受'概率云'")
    y = 1.5
    card = add_card(slide, ML, y, CONTENT_W, 1.5, LIGHT_GRAY)
    tf = card.text_frame; tf.word_wrap = True
    r = tf.paragraphs[0].add_run(); r.text = "1926 薛定谔方程提出的新描述方式"
    set_font(r, CN_TITLE, Pt(14), True, DARK_BLUE)
    add_paragraph(tf, "", CN_BODY, Pt(4))
    add_paragraph(tf, "  • 波函数 ψ 描述电子的运动状态（概率振幅）",
                  CN_BODY, Pt(13), color=DARK_GRAY, space_before=2)
    add_paragraph(tf, "  • |ψ|² 表示电子在空间某点出现的概率密度",
                  CN_BODY, Pt(13), True, DARK_BLUE, space_before=2)
    add_paragraph(tf, "  • 电子云 = |ψ|² 的可视化——黑点密集处概率大，稀疏处概率小",
                  CN_BODY, Pt(13), color=DARK_GRAY, space_before=2)
    y += 1.8
    add_textbox(slide, ML, y, 10, 0.5, "🗣️ 课堂原话：\"电子云不是云，是概率的散布图。这是从经典到量子的认知升级。\"",
                CN_BODY, Pt(13), False, DARK_GRAY, italic=True)
    # Three key takeaways
    y += 0.8
    takeaways = [
        ("🎯", "位置测不准", "无法同时知道位置和动量"),
        ("📊", "概率可确定", "但概率分布是确定的"),
        ("🧠", "统计规律", "用统计方法考察微观粒子行为"),
    ]
    for i, (icon, title, desc) in enumerate(takeaways):
        x = ML + i * 4.0
        c = add_card(slide, x, y, 3.5, 1.0, LIGHT_BLUE if i == 1 else LIGHT_GRAY)
        c_tf = c.text_frame; c_tf.word_wrap = True
        r = c_tf.paragraphs[0].add_run(); r.text = f"{icon} {title}"
        set_font(r, CN_TITLE, Pt(13), True, DARK_BLUE)
        add_paragraph(c_tf, desc, CN_BODY, Pt(11), color=DARK_GRAY, space_before=3)
    total += 1

    # ================================================================
    # SECTION 3: 原子结构的量子力学描述 (Slides 14-23)
    # ================================================================

    # --- Divider ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_divider(slide, "三", "原子结构的量子力学描述",
                        "薛定谔方程 · 四个量子数 · 概率密度与电子云")
    total += 1

    # --- Slide: 薛定谔方程 ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "薛定谔方程——量子力学的核心方程", "1926 · 以波函数 ψ 描述电子运动状态")
    y = 1.5
    # Equation card
    card = add_card(slide, ML, y, 11.0, 2.0, LIGHT_BLUE)
    tf = card.text_frame; tf.word_wrap = True
    r = tf.paragraphs[0].add_run(); r.text = "📐 薛定谔方程（直角坐标形式）"
    set_font(r, CN_TITLE, Pt(14), True, DARK_BLUE)
    add_paragraph(tf, "", CN_BODY, Pt(6))
    add_paragraph(tf, "  ∂²ψ/∂x² + ∂²ψ/∂y² + ∂²ψ/∂z² + (8π²m/h²)(E-V)ψ = 0",
                  CN_BODY, Pt(16), True, DARK_GRAY, PP_ALIGN.CENTER)
    add_paragraph(tf, "", CN_BODY, Pt(6))
    add_paragraph(tf, "  h: Planck常数   m: 电子质量   E: 总能量   V: 势能   ψ: 波函数（概率振幅）",
                  CN_BODY, Pt(11), False, GRAY, PP_ALIGN.CENTER)
    y += 2.5
    add_textbox(slide, ML, y, 10, 0.5, "💡 波函数的物理意义：对一个质量为 m、在势能为 V 的势场中运动的电子，",
                CN_BODY, Pt(13), False, DARK_GRAY)
    y += 0.4
    add_textbox(slide, ML + 0.3, y, 10, 0.5, "方程合理的解 ψ 表示电子运动的某一稳定状态，对应的 E 为该状态的电子能量",
                CN_BODY, Pt(13), False, DARK_GRAY)
    # Bottom: 坐标变换
    y += 0.6
    card = add_card(slide, ML, y, CONTENT_W, 1.2, LIGHT_GRAY)
    tf = card.text_frame; tf.word_wrap = True
    r = tf.paragraphs[0].add_run(); r.text = "由于原子具有球对称性 → 从直角坐标(x,y,z)变换到球极坐标(r,θ,φ)"
    set_font(r, CN_BODY, Pt(12), True, DARK_GRAY)
    add_paragraph(tf, "  x = r·sinθ·cosφ    y = r·sinθ·sinφ    z = r·cosθ",
                  CN_BODY, Pt(12), color=DARK_BLUE, space_before=4)
    total += 1

    # --- Slide: 三个量子数的自然出现 ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "变量分离与三个量子数的自然出现", "量子数不是人为假定的——是边界条件的自然推论")
    y = 1.3
    add_textbox(slide, ML, y, 10, 0.5, "通过变量分离法，将 ψ 分解为径向部分和角度部分：",
                CN_BODY, Pt(14), False, DARK_GRAY)
    y += 0.5
    card = add_card(slide, ML, y, 11.0, 0.7, LIGHT_AMBER)
    tf = card.text_frame; tf.word_wrap = True
    r = tf.paragraphs[0].add_run(); r.text = "    ψ_{n,l,m}(r,θ,φ) = R_{n,l}(r) · Y_{l,m}(θ,φ)"
    set_font(r, CN_BODY, Pt(18), True, DARK_BLUE, italic=True)
    y += 1.2
    # Three quantum numbers table
    q_data = [
        ("主量子数 n", "径向方程边界条件", "1, 2, 3, ...", "电子层、主要能量"),
        ("角量子数 l", "角度方程边界条件", "0, 1, ..., n-1", "轨道形状、能级分裂"),
        ("磁量子数 m_l", "角度方程边界条件", "-l, ..., 0, ..., +l", "轨道空间取向"),
    ]
    for i, (name, source, values, meaning) in enumerate(q_data):
        bg_c = LIGHT_GRAY if i % 2 == 0 else WHITE
        row_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            Inches(ML), Inches(y), Inches(CONTENT_W), Inches(0.55))
        row_bg.fill.solid(); row_bg.fill.fore_color.rgb = bg_c; row_bg.line.fill.background()
        add_textbox(slide, ML+0.2, y+0.05, 2.0, 0.4, name, CN_TITLE, Pt(13), True, DARK_BLUE)
        add_textbox(slide, ML+2.4, y+0.05, 3.0, 0.4, source, CN_BODY, Pt(11), False, GRAY)
        add_textbox(slide, ML+5.5, y+0.05, 2.5, 0.4, values, CN_BODY, Pt(10), True, DARK_GRAY)
        add_textbox(slide, ML+8.2, y+0.05, 3.0, 0.4, meaning, CN_BODY, Pt(11), False, DARK_GRAY)
        y += 0.55
    y += 0.2
    add_textbox(slide, ML, y, 10, 0.4, "🧠 关键认识：对应一组合理的 n,l,m 取值只有一个确定的波函数，每一个波函数称为一个原子轨道",
                CN_BODY, Pt(12), False, DARK_GRAY)
    total += 1

    # --- Slide: 主量子数 n ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "主量子数 n——电子层", "决定电子层和主要能量")
    y = 1.3
    # n-l table
    add_textbox(slide, ML, y, 4, 0.4, "n 取值与能层符号", CN_TITLE, Pt(14), True, DARK_BLUE)
    y += 0.5
    headers = ["n", "1", "2", "3", "4", "5", "6", "7"]
    symbol = ["K", "L", "M", "N", "O", "P", "Q"]
    # Header row
    for j, h in enumerate(headers):
        x = ML + j * 0.9
        tb = add_textbox(slide, x, y, 0.8, 0.35, h, CN_TITLE, Pt(12), True, DARK_BLUE, PP_ALIGN.CENTER)
    y += 0.35
    for j, s in enumerate(symbol):
        x = ML + (j+1) * 0.9
        bg_c = LIGHT_GRAY if j % 2 == 0 else WHITE
        row_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x-0.05), Inches(y), Inches(0.85), Inches(0.35))
        row_bg.fill.solid(); row_bg.fill.fore_color.rgb = bg_c; row_bg.line.fill.background()
        add_textbox(slide, x, y+0.02, 0.8, 0.3, s, CN_BODY, Pt(11), True, DARK_GRAY, PP_ALIGN.CENTER)
    # Key points right
    y_key = 1.3
    add_textbox(slide, 7.5, y_key, 5, 0.4, "📌 关键性质", CN_TITLE, Pt(14), True, DARK_BLUE)
    y_key += 0.5
    points = [
        "n 越大，电子离核越远，能量越高",
        "单电子体系：Eₙ = -13.6Z²/n² eV",
        "n 越大，原子半径越大",
        "轨道总数 = n²，最大电子数 = 2n²",
    ]
    for pt in points:
        add_textbox(slide, 7.5, y_key, 5, 0.35, f"  • {pt}", CN_BODY, Pt(12), False, DARK_GRAY)
        y_key += 0.35
    total += 1

    # --- Slide: 角量子数l ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "角量子数 l——轨道形状", "s球/p哑铃/d花瓣/f复杂")
    y = 1.3
    # l-table
    l_data = [
        (0, "s", "球形", "—"),
        (1, "p", "哑铃形", "3 种取向：px, py, pz"),
        (2, "d", "花瓣形", "5 种取向：dxy, dxz, dyz, dx²-y², dz²"),
        (3, "f", "复杂", "7 种取向"),
    ]
    add_textbox(slide, ML, y, 4, 0.4, "l 取值与轨道形状", CN_TITLE, Pt(14), True, DARK_BLUE)
    y += 0.5
    # Header
    for j, h in enumerate(["l值", "能级符号", "形状", "说明"]):
        xs = [ML, ML+1.2, ML+2.8, ML+4.5]
        add_textbox(slide, xs[j], y, 1.5, 0.35, h, CN_TITLE, Pt(11), True, DARK_BLUE)
    y += 0.35
    for i, (l, sym, shape, note) in enumerate(l_data):
        bg_c = LIGHT_GRAY if i % 2 == 0 else WHITE
        row_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            Inches(ML), Inches(y), Inches(5.8), Inches(0.35))
        row_bg.fill.solid(); row_bg.fill.fore_color.rgb = bg_c; row_bg.line.fill.background()
        add_textbox(slide, ML+0.1, y+0.02, 1.1, 0.3, str(l), CN_BODY, Pt(11), True, DARK_BLUE, PP_ALIGN.CENTER)
        add_textbox(slide, ML+1.2, y+0.02, 1.5, 0.3, sym, CN_BODY, Pt(11), True, DARK_GRAY, PP_ALIGN.CENTER)
        add_textbox(slide, ML+2.8, y+0.02, 1.5, 0.3, shape, CN_BODY, Pt(11), False, DARK_GRAY)
        add_textbox(slide, ML+4.5, y+0.02, 2.0, 0.3, note, CN_BODY, Pt(9), False, GRAY)
        y += 0.35
    # Right: energy ordering
    y_right = 1.3
    add_textbox(slide, 7.0, y_right, 5, 0.4, "📌 多电子原子能量顺序", CN_TITLE, Pt(14), True, DARK_BLUE)
    y_right += 0.5
    add_textbox(slide, 7.0, y_right, 5.5, 0.8,
                "当 n 相同时，l 越大 → 轨道能量越高：\nE(ns) < E(np) < E(nd) < E(nf)\n\n当 l 相同，n 越大 → 轨道能量越高：\nE(1s) < E(2s) < E(3s) < E(4s)",
                CN_BODY, Pt(12), False, DARK_GRAY)
    total += 1

    # --- Slide: 磁量子数m + 自旋量子数ms ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "磁量子数 m + 自旋量子数 m_s", "轨道空间取向 · 电子自旋")
    # Left: m magnetic
    y = 1.3
    add_textbox(slide, ML, y, 5, 0.4, "🧲 磁量子数 m", CN_TITLE, Pt(14), True, DARK_BLUE)
    y += 0.5
    add_textbox(slide, ML, y, 5.5, 0.35, "• 决定原子轨道的空间取向", CN_BODY, Pt(12), False, DARK_GRAY)
    y += 0.35
    add_textbox(slide, ML, y, 5.5, 0.35, "• 取值范围：-l, ..., 0, ..., +l（共 2l+1 个取值）", CN_BODY, Pt(12), False, DARK_GRAY)
    y += 0.35
    add_textbox(slide, ML, y, 5.5, 0.35, "• 轨道能量与 m 无关 → n,l相同 → 等价轨道（简并）", CN_BODY, Pt(12), False, DARK_GRAY)
    y += 0.5
    # Example: p orbitals
    add_textbox(slide, ML, y, 5, 0.35, "例：l=1(p轨道) → m=-1,0,+1 → px, py, pz 三个轨道", CN_BODY, Pt(11), True, DARK_BLUE)
    y += 0.4
    add_textbox(slide, ML, y, 5, 0.35, "例：l=2(d轨道) → m=-2,-1,0,+1,+2 → 5个d轨道", CN_BODY, Pt(11), True, DARK_BLUE)
    # Right: spin
    y2 = 1.3
    add_textbox(slide, 7.0, y2, 5, 0.4, "🔄 自旋量子数 m_s", CN_TITLE, Pt(14), True, DARK_BLUE)
    y2 += 0.5
    add_textbox(slide, 7.0, y2, 5.5, 0.35, "• 描述电子绕轴旋转的状态", CN_BODY, Pt(12), False, DARK_GRAY)
    y2 += 0.35
    add_textbox(slide, 7.0, y2, 5.5, 0.35, "• 取值：+1/2（↑）和 -1/2（↓）", CN_BODY, Pt(12), False, DARK_GRAY)
    y2 += 0.35
    add_textbox(slide, 7.0, y2, 5.5, 0.35, "• 自旋运动使电子具有类似微磁体的行为", CN_BODY, Pt(12), False, DARK_GRAY)
    y2 += 0.5
    add_textbox(slide, 7.0, y2, 5.5, 0.7,
                "💡 用四个量子数可以完整描述一个电子的运动状态：\nn(layer) · l(shape) · m(orientation) · ms(spin)",
                CN_BODY, Pt(11), True, DARK_BLUE)
    total += 1

    # --- Slide: 量子数汇总表 ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "量子数汇总表", "四个量子数共同决定一个电子的运动状态")
    # Build table
    table_data = [
        ["名称", "符号", "取值范围", "物理意义"],
        ["主量子数", "n", "1, 2, 3, ...", "电子层，主要能量"],
        ["角量子数", "l", "0, 1, ..., n-1", "轨道形状，能级分裂"],
        ["磁量子数", "m", "-l, ..., 0, ..., +l", "空间取向"],
        ["自旋量子数", "ms", "+1/2, -1/2", "自旋状态"],
    ]
    y = 1.4
    # Simple table-like layout
    col_x = [ML, ML+1.5, ML+3.5, ML+6.5]
    col_w = [1.4, 1.8, 2.8, 4.5]
    for i, row in enumerate(table_data):
        bg_c = DARK_BLUE if i == 0 else (LIGHT_GRAY if i % 2 == 1 else WHITE)
        row_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            Inches(ML), Inches(y), Inches(CONTENT_W), Inches(0.5))
        row_bg.fill.solid(); row_bg.fill.fore_color.rgb = bg_c; row_bg.line.fill.background()
        for j, cell in enumerate(row):
            is_header = (i == 0)
            font_cn = CN_TITLE if is_header else CN_BODY
            c = WHITE if is_header else DARK_GRAY
            add_textbox(slide, col_x[j], y+0.05, col_w[j], 0.4, cell,
                        font_cn, Pt(12), is_header, c, PP_ALIGN.CENTER)
        y += 0.5
    # Bottom summary
    y += 0.2
    total += 1

    # --- Slide: ⚠️ 量子数易错点 ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "⚠️ 易错点：量子数组合判断", "必须同时满足 n,l,m,ms 的取值范围！")
    y = 1.3
    add_textbox(slide, ML, y, 10, 0.4, "判断以下量子数组是否能存在？为什么？（选自原讲义习题1）",
                CN_BODY, Pt(14), True, DARK_GRAY)
    y += 0.6
    # Examples
    examples = [
        ("(1) n=1, l=1, m=1, ms=-1", "❌ l 不能 ≥ n（l max = n-1 = 0）"),
        ("(2) n=3, l=1, m=2, ms=+1/2", "❌ m 超出范围（l=1 → m = -1,0,+1）"),
        ("(3) n=3, l=2, m=1, ms=-1/2", "✅ 合法：3d 轨道的一个电子"),
        ("(4) n=2, l=0, m=0, ms=0", "❌ ms 不能为 0（只能 +1/2 或 -1/2）"),
    ]
    for i, (q, a) in enumerate(examples):
        bg_c = LIGHT_GREEN if a.startswith("✅") else LIGHT_RED
        card = add_card(slide, ML, y, CONTENT_W, 0.55, bg_c)
        tf = card.text_frame; tf.word_wrap = True
        r = tf.paragraphs[0].add_run(); r.text = f"  {q}"
        set_font(r, CN_BODY, Pt(12), True, DARK_GRAY)
        p = add_paragraph(tf, f"  {a}", CN_BODY, Pt(11), True,
                          GREEN if "✅" in a else RED, space_before=1)
        y += 0.6
    total += 1

    # --- Slide: 概率密度与电子云 ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "概率密度与电子云", "|ψ|² = 概率密度")
    y = 1.3
    add_textbox(slide, ML, y, 6, 0.4, "|ψ|² 的物理意义：概率密度（单位体积内电子出现概率）",
                CN_BODY, Pt(14), True, DARK_BLUE)
    y += 0.5
    concepts = [
        ("径向密度分布", "R²(r) 对 r 作图 — 核附近电子云密度最大，随 r 增大递减\ns 态在核附近密度最大；p/d 态核附近密度接近零"),
        ("径向分布函数", "D(r) = r²R²(r) — 单位厚度球壳内电子分布概率\n1s 的 D 函数极大值在 Bohr 半径 a₀=52.9pm 处"),
        ("峰的规律", "• D 函数峰个数 = n - l\n• 主峰随 n 增加离核越来越远\n• l 越大，峰数目越少，主峰离核越近"),
    ]
    for title, body in concepts:
        card = add_card(slide, ML, y, CONTENT_W, 1.1, LIGHT_GRAY)
        tf = card.text_frame; tf.word_wrap = True
        r = tf.paragraphs[0].add_run(); r.text = f"  {title}"
        set_font(r, CN_TITLE, Pt(12), True, DARK_BLUE)
        add_paragraph(tf, f"  {body}", CN_BODY, Pt(10), color=DARK_GRAY, space_before=2)
        y += 1.2
    total += 1

    # --- Slide: 原子轨道与电子云角度分布 ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "原子轨道与电子云的空间图象", "波函数角度分布 vs 电子云角度分布")
    # Comparison table
    y = 1.3
    add_textbox(slide, ML, y, 6, 0.4, "波函数(ψ) 与 电子云(ψ²) 对比",
                CN_TITLE, Pt(14), True, DARK_BLUE)
    y += 0.5
    comp_data = [
        ["对比项", "ψ 波函数", "ψ² 电子云"],
        ["定义", "描写电子运动状态的数学表达式", "核外空间某处电子分布概率"],
        ["形态", "角度分布图略\"胖\"", "比波函数图略\"瘦\""],
        ["符号", "有正负(+/−)", "无符号（均为正）"],
        ["应用", "分析化学键的形成", "讨论键的空间构型"],
    ]
    col_x2 = [ML, ML+3.5, ML+8.0]
    col_w2 = [3.3, 4.3, 4.3]
    for i, row in enumerate(comp_data):
        bg_c = DARK_BLUE if i == 0 else (LIGHT_GRAY if i % 2 == 1 else WHITE)
        row_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            Inches(ML), Inches(y), Inches(CONTENT_W), Inches(0.45))
        row_bg.fill.solid(); row_bg.fill.fore_color.rgb = bg_c; row_bg.line.fill.background()
        for j, cell in enumerate(row):
            is_header = (i == 0)
            c = WHITE if is_header else DARK_GRAY
            f = CN_TITLE if is_header else CN_BODY
            add_textbox(slide, col_x2[j], y+0.05, col_w2[j], 0.35, cell,
                        f, Pt(11), is_header, c, PP_ALIGN.CENTER)
        y += 0.45
    total += 1

    # ================================================================
    # SECTION 4: 原子核外电子排布 (Slides 24-31)
    # ================================================================

    # --- Divider ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_divider(slide, "四", "原子核外电子排布",
                        "构造原理 · 屏蔽与钻穿 · Slater 规则 · 三原则")
    total += 1

    # --- Slide: 构造原理与近似能级图 ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "构造原理与近似能级图", "鲍林光谱实验总结 → 能级组对应周期")
    y = 1.3
    add_textbox(slide, ML, y, 10, 0.4, "电子填充顺序按能量从低到高排列：",
                CN_BODY, Pt(14), True, DARK_GRAY)
    y += 0.5
    # Pauling order
    order = "1s  2s 2p  3s 3p  4s 3d 4p  5s 4d 5p  6s 4f 5d 6p  7s 5f 6d 7p"
    card = add_card(slide, ML, y, CONTENT_W, 0.7, LIGHT_AMBER)
    tf = card.text_frame; tf.word_wrap = True
    r = tf.paragraphs[0].add_run(); r.text = f"  {order}"
    set_font(r, CN_BODY, Pt(14), True, DARK_BLUE, italic=True)
    y += 1.0
    # 徐光宪规则
    add_textbox(slide, ML, y, 10, 0.4, "徐光宪规则：按 n+0.7l 值排序（值越大能量越高）",
                CN_TITLE, Pt(13), True, DARK_BLUE)
    y += 0.45
    examples_xu = [
        "例：第七能级组 — 7s(7.0) → 5f(6.4) → 6d(6.7) → 7p(7.7)，同属第7周期",
        "n+0.7l 整数部分相同 → 同属一个能级组 → 对应一个周期",
    ]
    for ex in examples_xu:
        add_textbox(slide, ML + 0.3, y, 10, 0.35, f"  • {ex}", CN_BODY, Pt(11), False, DARK_GRAY)
        y += 0.35
    y += 0.2
    add_callout_box(slide, ML, y, CONTENT_W, 0.9, "🧠", "教学洞察：能级交错——填充顺序不是 n 的自然顺序",
                    "4s 先于 3d → 钻穿效应导致 4s 能量低于 3d\nE(4s) < E(3d) < E(4p) — 详见 §钻穿效应",
                    LIGHT_AMBER)
    total += 1

    # --- Slide: 屏蔽效应 ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "屏蔽效应", "其他电子抵消部分核电荷 → 有效核电荷 Z* < Z")
    y = 1.3
    add_textbox(slide, ML, y, 10, 0.4, "概念：电子 i 受其他电子的排斥，抵消了部分核电荷的吸引",
                CN_BODY, Pt(14), True, DARK_GRAY)
    y += 0.5
    card = add_card(slide, ML, y, 6.0, 2.0, LIGHT_BLUE)
    tf = card.text_frame; tf.word_wrap = True
    r = tf.paragraphs[0].add_run(); r.text = "📐 有效核电荷"
    set_font(r, CN_TITLE, Pt(14), True, DARK_BLUE)
    add_paragraph(tf, "", CN_BODY, Pt(6))
    add_paragraph(tf, "    Z* = Z - σ", CN_BODY, Pt(22), True, DARK_GRAY, PP_ALIGN.CENTER)
    add_paragraph(tf, "", CN_BODY, Pt(4))
    add_paragraph(tf, "    Z: 核电荷数    σ: 屏蔽常数    Z*: 有效核电荷",
                  CN_BODY, Pt(11), False, GRAY, PP_ALIGN.CENTER)
    add_paragraph(tf, "", CN_BODY, Pt(4))
    add_paragraph(tf, "    屏蔽越大 → σ 越大 → Z* 越小 → 电子能量越高",
                  CN_BODY, Pt(11), False, DARK_GRAY, PP_ALIGN.CENTER)
    y_right = 1.3
    add_textbox(slide, 7.2, y_right, 5.5, 0.4, "🔑 屏蔽效应的影响", CN_TITLE, Pt(14), True, DARK_BLUE)
    y_right += 0.5
    points = [
        "同层电子间屏蔽作用较小（σ≈0.35）",
        "内层对外层屏蔽作用接近完全（σ≈0.85-1.00）",
        "外层对内层无屏蔽作用（σ=0）",
        "屏蔽效应是能级交错的重要原因之一",
    ]
    for pt in points:
        add_textbox(slide, 7.2, y_right, 5.5, 0.35, f"  • {pt}", CN_BODY, Pt(11), False, DARK_GRAY)
        y_right += 0.35
    total += 1

    # --- Slide: Slater规则 ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "Slater 经验规则——屏蔽常数的定量计算", "🌟了解性：会算 Na 和 Ti 即可，不要求复杂体系")
    y = 1.3
    add_textbox(slide, ML, y, 10, 0.4, "Slater 规则将原子轨道分为若干组：",
                CN_BODY, Pt(14), False, DARK_GRAY)
    y += 0.45
    add_textbox(slide, ML + 0.3, y, 10, 0.4, "(1s) (2s2p) (3s3p) (3d) (4s4p) (4d) (4f) (5s5p) (5d) ...",
                CN_BODY, Pt(12), True, DARK_BLUE, italic=True)
    y += 0.5
    rules = [
        "外层电子对内层无屏蔽（σ=0）",
        "同组电子间 σ=0.35（1s 组内 σ=0.30）",
        "对 ns/np 电子：(n-1)组 σ=0.85；(n-2)及更内 σ=1.00",
        "对 nd/nf 电子：左侧所有组 σ=1.00",
    ]
    for rule in rules:
        add_textbox(slide, ML + 0.3, y, 10, 0.35, f"  {rule}", CN_BODY, Pt(12), False, DARK_GRAY)
        y += 0.35
    y += 0.2
    card = add_card(slide, ML, y, CONTENT_W, 1.0, LIGHT_GREEN)
    tf = card.text_frame; tf.word_wrap = True
    r = tf.paragraphs[0].add_run(); r.text = "💡 第一轮要求：Na（1s²2s²2p⁶3s¹）和 Ti（1s²...3d²4s²）的完整计算示例"
    set_font(r, CN_BODY, Pt(12), True, DARK_GRAY)
    add_paragraph(tf, "  Slater 规则的目的是让学生定量感受\"为什么 4s 能量低于 3d\"——不是背规则本身。",
                  CN_BODY, Pt(11), False, GRAY, space_before=3)
    total += 1

    # --- Slide: 钻穿效应 ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "钻穿效应——4s 为什么先于 3d？", "电子钻入内层空间 → 受核引力更强 → 能量降低")
    y = 1.3
    add_textbox(slide, ML, y, 10, 0.4, "概念：电子进入原子内部空间，受到核的较强吸引作用",
                CN_BODY, Pt(14), True, DARK_GRAY)
    y += 0.5
    # Mechanism
    card = add_card(slide, ML, y, 6.0, 2.2, LIGHT_GRAY)
    tf = card.text_frame; tf.word_wrap = True
    r = tf.paragraphs[0].add_run(); r.text = "🔬 物理图像"
    set_font(r, CN_TITLE, Pt(13), True, DARK_BLUE)
    add_paragraph(tf, "", CN_BODY, Pt(4))
    add_paragraph(tf, "  • 4s 电子云有小峰钻到离核很近的地方",
                  CN_BODY, Pt(12), color=DARK_GRAY, space_before=3)
    add_paragraph(tf, "  • 3d 电子云被 3s/3p 挡住，\"感受不到核的引力\"",
                  CN_BODY, Pt(12), color=DARK_GRAY, space_before=2)
    add_paragraph(tf, "  • 电子钻得越深 → 受核吸引力越强 → 能量越低",
                  CN_BODY, Pt(12), color=DARK_GRAY, space_before=2)
    add_paragraph(tf, "  • n 相同时：钻穿效应 ns > np > nd > nf",
                  CN_BODY, Pt(12), True, DARK_BLUE, space_before=2)
    # Right: analogy
    add_textbox(slide, 7.2, y, 5.5, 0.4, "🗣️ 类比：\"3d 电子被 3s/3p 电子云挡住——", CN_BODY, Pt(12), False, DARK_GRAY, italic=True)
    add_textbox(slide, 7.2, y+0.4, 5.5, 0.4, "像躲在盾牌后面感受不到核的引力\"", CN_BODY, Pt(12), False, DARK_GRAY, italic=True)
    # Key comparison
    y2 = 2.8
    card2 = add_card(slide, 7.2, y2, 5.5, 0.7, LIGHT_AMBER)
    tf2 = card2.text_frame; tf2.word_wrap = True
    r = tf2.paragraphs[0].add_run(); r.text = "🧠 结论：E(4s) < E(3d) < E(4p)"
    set_font(r, CN_TITLE, Pt(13), True, DARK_BLUE)
    add_paragraph(tf2, "原因：4s 钻穿效应 > 3d → 4s 能量降低更多",
                  CN_BODY, Pt(11), color=DARK_GRAY, space_before=2)
    total += 1

    # --- Slide: 电子排布三原则 ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "电子排布三原则", "能量最低 · 泡利不相容 · 洪特规则")
    y = 1.3
    principles = [
        ("⚡ 能量最低原则", "电子在原子轨道中排布，要使整个原子系统能量最低\n按构造原理顺序填充：1s→2s→2p→3s→3p→4s→3d→4p→...", DARK_BLUE, LIGHT_BLUE),
        ("🚫 泡利不相容原理", "每个原子轨道最多容纳 2 个自旋相反的电子\n同一原子中不可能有 4 个量子数完全相同的两个电子\n→ 一个轨道中不存在自旋相同的两个电子", BRIGHT_BLUE, LIGHT_GRAY),
        ("🪑 洪特规则", "电子在简并轨道(能量相同)上排布时\n尽可能以自旋平行的方向，单独分占不同轨道\n特例：半满、全满、全空时特别稳定", AMBER, LIGHT_AMBER),
    ]
    for i, (title, desc, clr, bg) in enumerate(principles):
        card = add_card(slide, ML, y, CONTENT_W, 1.3, bg)
        tf = card.text_frame; tf.word_wrap = True
        r = tf.paragraphs[0].add_run(); r.text = title
        set_font(r, CN_TITLE, Pt(14), True, clr)
        add_paragraph(tf, f"  {desc}", CN_BODY, Pt(11), color=DARK_GRAY, space_before=3)
        y += 1.4
    # Bottom: 先占座后配对
    y += 0.1
    add_callout_box(slide, ML, y, CONTENT_W, 0.7, "🗣️", "课堂原话：\"先占座后配对\"——公交车上先坐空座，没空座了才挤（洪特规则）",
                    "类比来源：质心L2 · 原话",
                    LIGHT_GREEN)
    total += 1

    # --- Slide: Cr/Cu特例 ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "Cr/Cu 特例——半满全满稳定化", "\"反常\"排布比\"正常\"排布更稳定")
    y = 1.3
    # Comparison
    left_data = [
        ("\n元素", "预期排布", "实际排布", "原因"),
        ("Cr (Z=24)", "[Ar]4s²3d⁴", "[Ar]4s¹3d⁵ ✅", "3d 半满稳定"),
        ("Cu (Z=29)", "[Ar]4s²3d⁹", "[Ar]4s¹3d¹⁰ ✅", "3d 全满稳定"),
        ("Mo (Z=42)", "[Kr]5s²4d⁴", "[Kr]5s¹4d⁵ ✅", "4d 半满（Cr 族）"),
        ("Ag (Z=47)", "[Kr]5s²4d⁹", "[Kr]5s¹4d¹⁰ ✅", "4d 全满（Cu 族）"),
    ]
    col_x3 = [ML, ML+1.5, ML+4.0, ML+6.5]
    col_w3 = [1.4, 2.3, 2.3, 3.0]
    for i, row in enumerate(left_data):
        bg_c = DARK_BLUE if i == 0 else (LIGHT_GRAY if i % 2 == 1 else WHITE)
        row_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            Inches(ML), Inches(y), Inches(8.0), Inches(0.5))
        row_bg.fill.solid(); row_bg.fill.fore_color.rgb = bg_c; row_bg.line.fill.background()
        for j, cell in enumerate(row):
            is_h = (i == 0)
            c = WHITE if is_h else DARK_GRAY
            f = CN_TITLE if is_h else CN_BODY
            add_textbox(slide, col_x3[j], y+0.05, col_w3[j], 0.4, cell,
                        f, Pt(11), is_h, c, PP_ALIGN.CENTER)
        y += 0.5
    # Right rules
    y_r = 1.3
    add_textbox(slide, 8.2, y_r, 4.5, 0.4, "📌 记忆口诀", CN_TITLE, Pt(14), True, DARK_BLUE)
    y_r += 0.5
    add_textbox(slide, 8.2, y_r, 4.5, 0.8,
                "\"半满全满更稳定\nCr 和 Cu 要记住\"\n\n仅适用于 Cr 族和 Cu 族\n其余过渡金属不适用！",
                CN_BODY, Pt(12), False, DARK_GRAY)
    total += 1

    # --- Slide: 填充vs失电子 ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "🧠 关键认知冲突：填充顺序 ≠ 失电子顺序", "\"入住顺序 ≠ 退房顺序\"——4s 是酒店大堂")
    y = 1.3
    # Analogy card
    card = add_card(slide, ML, y, CONTENT_W, 2.2, LIGHT_AMBER)
    tf = card.text_frame; tf.word_wrap = True
    r = tf.paragraphs[0].add_run(); r.text = "🏨 酒店大堂类比"
    set_font(r, CN_TITLE, Pt(16), True, DARK_BLUE)
    add_paragraph(tf, "", CN_BODY, Pt(4))
    add_paragraph(tf, "  入住时：先到大堂（4s），再到房间（3d）  →  填充顺序：4s 先于 3d",
                  CN_BODY, Pt(13), True, DARK_GRAY, space_before=4)
    add_paragraph(tf, "  退房时：从大堂走（4s），房间还在（3d） →  失电子顺序：4s 先于 3d",
                  CN_BODY, Pt(13), True, DARK_GRAY, space_before=3)
    add_paragraph(tf, "", CN_BODY, Pt(4))
    add_paragraph(tf, "  关键：4s 和 3d 的能量本来就接近，填充时 4s 略低→先填",
                  CN_BODY, Pt(12), False, DARK_GRAY, space_before=3)
    add_paragraph(tf, "  但写排布式时，按 n 顺序写：1s²2s²2p⁶3s²3p⁶3d...4s...",
                  CN_BODY, Pt(12), True, DARK_BLUE, space_before=2)
    # Error rate
    y += 2.5
    card2 = add_card(slide, ML, y, CONTENT_W, 0.8, LIGHT_RED)
    tf2 = card2.text_frame; tf2.word_wrap = True
    r = tf2.paragraphs[0].add_run(); r.text = "⚠️ 高频错误（发生率 ~25%）：认为离子失电子也按填充顺序"
    set_font(r, CN_BODY, Pt(12), True, RED)
    add_paragraph(tf2, "  正解：Fe²⁺ = [Ar]3d⁶（先失 4s²，不是 3d⁶）",
                  CN_BODY, Pt(12), True, DARK_GRAY, space_before=3)
    total += 1

    # --- Slide: 过渡金属排布 ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "过渡金属电子排布（Sc-Zn）", "第一过渡系 21-30 号元素的价电子排布")
    y = 1.3
    # Create transition metal table
    elements = [
        ("Sc", "21", "[Ar]4s²3d¹"), ("Ti", "22", "[Ar]4s²3d²"),
        ("V", "23", "[Ar]4s²3d³"), ("Cr", "24", "[Ar]4s¹3d⁵"),
        ("Mn", "25", "[Ar]4s²3d⁵"), ("Fe", "26", "[Ar]4s²3d⁶"),
        ("Co", "27", "[Ar]4s²3d⁷"), ("Ni", "28", "[Ar]4s²3d⁸"),
        ("Cu", "29", "[Ar]4s¹3d¹⁰"), ("Zn", "30", "[Ar]4s²3d¹⁰"),
    ]
    # Headers
    hdrs = ["Z", "21", "22", "23", "24", "25", "26", "27", "28", "29", "30"]
    el_sym = ["Sc", "Ti", "V", "Cr", "Mn", "Fe", "Co", "Ni", "Cu", "Zn"]
    el_cfg = ["[Ar]4s²3d¹", "[Ar]4s²3d²", "[Ar]4s²3d³", "[Ar]4s¹3d⁵", "[Ar]4s²3d⁵",
              "[Ar]4s²3d⁶", "[Ar]4s²3d⁷", "[Ar]4s²3d⁸", "[Ar]4s¹3d¹⁰", "[Ar]4s²3d¹⁰"]
    # Row 1: Element symbols
    y_row = y
    add_textbox(slide, ML, y_row, 1.0, 0.35, "元素", CN_TITLE, Pt(11), True, DARK_BLUE)
    for j, s in enumerate(el_sym):
        x = ML + 1.1 + j * 1.05
        bg_c = LIGHT_AMBER if s in ["Cr", "Cu"] else LIGHT_GRAY
        row_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            Inches(x-0.05), Inches(y_row), Inches(1.0), Inches(0.35))
        row_bg.fill.solid(); row_bg.fill.fore_color.rgb = bg_c; row_bg.line.fill.background()
        c = RED if s in ["Cr", "Cu"] else DARK_GRAY
        add_textbox(slide, x, y_row+0.02, 0.9, 0.3, s, CN_TITLE, Pt(11), True, c, PP_ALIGN.CENTER)
    # Row 2: Configurations
    y_row += 0.4
    add_textbox(slide, ML, y_row, 1.0, 0.35, "排布", CN_TITLE, Pt(11), True, DARK_BLUE)
    for j, cfg in enumerate(el_cfg):
        x = ML + 1.1 + j * 1.05
        add_textbox(slide, x, y_row+0.02, 1.0, 0.3, cfg,
                    CN_BODY, Pt(7), True, DARK_GRAY, PP_ALIGN.CENTER)
    # Notes
    y += 1.2
    add_textbox(slide, ML, y, 10, 0.4, "⚠️ 红色标注的 Cr 和 Cu 为特例——半满/全满稳定化",
                CN_BODY, Pt(12), True, RED)
    total += 1

    # ================================================================
    # SECTION 5: 核外电子排布（续）(Slides 32-34)
    # ================================================================

    # --- Divider ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_divider(slide, "五", "核外电子排布（续）· 周期律",
                        "离子排布 · 原子半径 · 电离能 · 电负性 · 反常点")
    total += 1

    # --- Slide: 离子电子排布 ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "离子电子排布——先失最外层", "失电子顺序：np → ns → (n-1)d")
    y = 1.3
    add_textbox(slide, ML, y, 10, 0.5, "规则：正离子先失电子从最外层开始，按 np → ns → (n-1)d 顺序",
                CN_BODY, Pt(14), True, DARK_BLUE)
    y += 0.6
    examples_ion = [
        ("Fe → Fe²⁺ → Fe³⁺", "Fe [Ar]4s²3d⁶ → Fe²⁺ [Ar]3d⁶（先失4s²）→ Fe³⁺ [Ar]3d⁵（再失1个3d）"),
        ("Cu → Cu⁺ → Cu²⁺", "Cu [Ar]4s¹3d¹⁰ → Cu⁺ [Ar]3d¹⁰（失4s¹）→ Cu²⁺ [Ar]3d⁹（再失1个3d）"),
        ("Mn²⁺", "Mn [Ar]4s²3d⁵ → Mn²⁺ [Ar]3d⁵（3d 半满，特别稳定）"),
        ("Cr³⁺", "Cr [Ar]4s¹3d⁵ → Cr³⁺ [Ar]3d³"),
    ]
    for title, body in examples_ion:
        card = add_card(slide, ML, y, CONTENT_W, 0.6, LIGHT_GRAY)
        tf = card.text_frame; tf.word_wrap = True
        r = tf.paragraphs[0].add_run(); r.text = f"  {title}"
        set_font(r, CN_TITLE, Pt(12), True, DARK_BLUE)
        add_paragraph(tf, f"  {body}", CN_BODY, Pt(10), color=DARK_GRAY, space_before=1)
        y += 0.65
    total += 1

    # --- Slide: 周期律趋势 ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "周期律三大参数趋势", "原子半径 · 电离能 · 电负性")
    y = 1.3
    # Three trends
    trends = [
        ("📏 原子半径", "同周期左→右：递减（Z*↑，电子云收缩）\n同主族上→下：递增（n↑为主）\n镧系收缩：La→Lu 缓慢递减"),
        ("⚡ 电离能 I", "同周期左→右：递增（Z*↑，电子难失去）\n同主族上→下：递减（n↑，半径↑）\n⚠️ 反常点：ⅡA>ⅢA，ⅤA>ⅥA"),
        ("🔗 电负性 χ", "同周期左→右：递增\n同主族上→下：递减\nF 最大(4.0)，Cs 最小(0.7)\n差值大→离子键；差值小→共价键"),
    ]
    for i, (title, desc, *_) in enumerate(trends):
        x = ML + i * 4.0
        card = add_card(slide, x, y, 3.7, 2.5, LIGHT_BLUE if i == 1 else LIGHT_GRAY)
        tf = card.text_frame; tf.word_wrap = True
        r = tf.paragraphs[0].add_run(); r.text = title
        set_font(r, CN_TITLE, Pt(14), True, DARK_BLUE)
        add_paragraph(tf, f"  {desc}", CN_BODY, Pt(11), color=DARK_GRAY, space_before=4)
    total += 1

    # --- Slide: 电离能反常 ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "⚠️ 电离能反常清单", "只背趋势不看反常 → 做题必错（发生率~30%）")
    y = 1.3
    add_textbox(slide, ML, y, 10, 0.5, "口诀：\"比较电离能，先看轨道类型，再看周期趋势\"",
                CN_BODY, Pt(15), True, DARK_BLUE)
    y += 0.6
    # Anomaly table
    anomalies = [
        ("B vs Be", "Be 1s²2s²（s 全满稳定）", "B 2s²2p¹（p 电子易于电离）", "Be I₁ > B I₁"),
        ("O vs N", "N 2s²2p³（p 半满稳定）", "O 2s²2p⁴（p⁴有两个电子同轨道排斥）", "N I₁ > O I₁"),
        ("Mg vs Al", "Mg 3s²（s 全满稳定）", "Al 3s²3p¹（p 电子易于电离）", "Mg I₁ > Al I₁"),
        ("P vs S", "P 3s²3p³（p 半满稳定）", "S 3s²3p⁴（p⁴有电子同轨道排斥）", "P I₁ > S I₁"),
    ]
    col_x4 = [ML, ML+2.5, ML+5.5, ML+8.5]
    col_w4 = [2.3, 2.8, 2.8, 2.2]
    # Header
    hdr_row = ["反常对", "更高I₁的原因", "更低I₁的原因", "结论"]
    for j, (h, x, w) in enumerate(zip(hdr_row, col_x4, col_w4)):
        add_textbox(slide, x, y, w, 0.35, h, CN_TITLE, Pt(10), True, WHITE)
    # Manually set white text on header by drawing a blue card
    hdr_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(ML), Inches(y), Inches(CONTENT_W), Inches(0.35))
    hdr_bg.fill.solid(); hdr_bg.fill.fore_color.rgb = DARK_BLUE; hdr_bg.line.fill.background()
    for j, (h, x, w) in enumerate(zip(hdr_row, col_x4, col_w4)):
        add_textbox(slide, x, y+0.02, w, 0.3, h, CN_TITLE, Pt(10), True, WHITE, PP_ALIGN.CENTER)
    y += 0.35
    for i, (pair, hi, lo, conclusion) in enumerate(anomalies):
        bg_c = LIGHT_GRAY if i % 2 == 0 else WHITE
        row_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            Inches(ML), Inches(y), Inches(CONTENT_W), Inches(0.45))
        row_bg.fill.solid(); row_bg.fill.fore_color.rgb = bg_c; row_bg.line.fill.background()
        add_textbox(slide, col_x4[0], y+0.05, col_w4[0], 0.35, pair,
                    CN_TITLE, Pt(11), True, DARK_BLUE, PP_ALIGN.CENTER)
        add_textbox(slide, col_x4[1], y+0.05, col_w4[1], 0.35, hi,
                    CN_BODY, Pt(9), False, DARK_GRAY)
        add_textbox(slide, col_x4[2], y+0.05, col_w4[2], 0.35, lo,
                    CN_BODY, Pt(9), False, DARK_GRAY)
        add_textbox(slide, col_x4[3], y+0.05, col_w4[3], 0.35, conclusion,
                    CN_BODY, Pt(10), True, RED, PP_ALIGN.CENTER)
        y += 0.45
    total += 1

    # ================================================================
    # SECTION 6: 原子核衰变 (Slides 35-37)
    # ================================================================

    # --- Divider ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_divider(slide, "六", "原子核衰变",
                        "α 衰变 · β 衰变 · γ 衰变 · 核反应方程式")
    total += 1

    # --- Slide: 三种衰变 ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "三种核衰变类型对比", "放射性核素自发释放射线 → 转变成另一种核")
    y = 1.3
    decays = [
        ("α 衰变", "⁴₂He (α粒子)", "质量数 -4\n核电荷 -2", "U → Th（核武器）"),
        ("β 衰变", "⁰₋₁e (β粒子/电子)", "质量数不变\n核电荷 +1", "¹⁴C → ¹⁴N（碳定年）"),
        ("γ 衰变", "γ 光子（高能电磁波）", "质量数不变\n核电荷不变", "激发态→基态退激"),
    ]
    col_x5 = [ML, ML+2.5, ML+5.0, ML+8.0]
    col_w5 = [2.3, 2.3, 2.8, 2.5]
    # Header
    hdr_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(ML), Inches(y), Inches(CONTENT_W), Inches(0.4))
    hdr_bg.fill.solid(); hdr_bg.fill.fore_color.rgb = DARK_BLUE; hdr_bg.line.fill.background()
    for j, (h, x, w) in enumerate(zip(["衰变类型", "放出粒子", "核变化", "实例"], col_x5, col_w5)):
        add_textbox(slide, x, y+0.05, w, 0.3, h, CN_TITLE, Pt(12), True, WHITE, PP_ALIGN.CENTER)
    y += 0.4
    for i, (name, particle, change, example) in enumerate(decays):
        bg_c = LIGHT_GRAY if i % 2 == 0 else WHITE
        row_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            Inches(ML), Inches(y), Inches(CONTENT_W), Inches(0.8))
        row_bg.fill.solid(); row_bg.fill.fore_color.rgb = bg_c; row_bg.line.fill.background()
        c = DARK_BLUE if i == 2 else (BRIGHT_BLUE if i == 1 else AMBER)
        add_textbox(slide, col_x5[0], y+0.15, col_w5[0], 0.5, name,
                    CN_TITLE, Pt(14), True, c, PP_ALIGN.CENTER)
        add_textbox(slide, col_x5[1], y+0.15, col_w5[1], 0.5, particle,
                    CN_BODY, Pt(11), False, DARK_GRAY, PP_ALIGN.CENTER)
        add_textbox(slide, col_x5[2], y+0.10, col_w5[2], 0.6, change,
                    CN_BODY, Pt(11), False, DARK_GRAY, PP_ALIGN.CENTER)
        add_textbox(slide, col_x5[3], y+0.15, col_w5[3], 0.5, example,
                    CN_BODY, Pt(10), italic=True, color=GRAY, alignment=PP_ALIGN.CENTER)
        y += 0.85
    total += 1

    # --- Slide: 核反应方程式 ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "核化学方程式书写规则", "原子核反应的两条守恒定律")
    y = 1.5
    rules_nuclear = [
        ("规则一", "质量数守恒：方程式两端总质量数相等"),
        ("规则二", "原子序数守恒：方程式两端总电荷数相等"),
        ("规则三", "α/β/γ 衰变常伴有 γ 射线放射"),
    ]
    for i, (title, desc) in enumerate(rules_nuclear):
        card = add_card(slide, ML, y, CONTENT_W, 0.6, LIGHT_BLUE if i < 2 else LIGHT_GRAY)
        tf = card.text_frame; tf.word_wrap = True
        r = tf.paragraphs[0].add_run(); r.text = f"  {title}"
        set_font(r, CN_TITLE, Pt(13), True, DARK_BLUE)
        add_paragraph(tf, f"  {desc}", CN_BODY, Pt(12), color=DARK_GRAY, space_before=2)
        y += 0.7
    # Example
    y += 0.2
    card = add_card(slide, ML, y, CONTENT_W, 0.7, LIGHT_AMBER)
    tf = card.text_frame; tf.word_wrap = True
    r = tf.paragraphs[0].add_run(); r.text = "  例：²³⁸U → ²³⁴Th + ⁴₂He"
    set_font(r, CN_BODY, Pt(16), True, DARK_GRAY, italic=True)
    add_paragraph(tf, "  质量数：238 = 234 + 4 ✓    核电荷：92 = 90 + 2 ✓",
                  CN_BODY, Pt(12), False, DARK_GRAY, space_before=3)
    total += 1

    # --- Slide: 半衰期 ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "半衰期", "放射性元素的原子核有半数发生衰变所需时间")
    y = 1.5
    add_textbox(slide, ML, y, 10, 0.5, "定义：放射性原子核数衰变到原来一半所需的时间（t₁/₂）",
                CN_BODY, Pt(14), True, DARK_GRAY)
    y += 0.6
    card = add_card(slide, ML, y, 6.0, 1.8, LIGHT_BLUE)
    tf = card.text_frame; tf.word_wrap = True
    r = tf.paragraphs[0].add_run(); r.text = "📐 衰变规律"
    set_font(r, CN_TITLE, Pt(14), True, DARK_BLUE)
    add_paragraph(tf, "", CN_BODY, Pt(4))
    add_paragraph(tf, "    N = N₀ · (1/2)^(t/t₁/₂)", CN_BODY, Pt(18), True, DARK_GRAY, PP_ALIGN.CENTER)
    add_paragraph(tf, "", CN_BODY, Pt(4))
    add_paragraph(tf, "    N₀: 初始原子数    N: 衰变后剩余数    t: 时间",
                  CN_BODY, Pt(11), False, GRAY, PP_ALIGN.CENTER)
    total += 1

    # ================================================================
    # TYPICAL EXAMPLES (Slides 38-40)
    # ================================================================

    # --- Divider ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_divider(slide, "", "典型例题", "精选竞赛真题 · 覆盖面：量子数 → 排布 → 核反应")
    total += 1

    # --- Slide: 例1-量子数 ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "例题·量子数与轨道", "2000年美国国家化学竞赛试题")
    y = 1.5
    add_textbox(slide, ML, y, 10, 0.8,
                "在多电子原子中，下列关于电子依次连续填入轨道顺序描述正确的是：\nA. 3s, 3p, 3d    B. 3d, 4s, 4p    C. 3d, 4p, 5s    D. 4p, 4d, 5s",
                CN_BODY, Pt(14), False, DARK_GRAY)
    y += 1.0
    card = add_card(slide, ML, y, CONTENT_W, 1.0, LIGHT_GREEN)
    tf = card.text_frame; tf.word_wrap = True
    r = tf.paragraphs[0].add_run(); r.text = "  ✅ 正确答案：B"
    set_font(r, CN_TITLE, Pt(16), True, GREEN)
    add_paragraph(tf, "  🔍 解析：填充顺序按能量从低到高，3d≈4s→4p，B 的 3d→4s→4p 最接近实际顺序",
                  CN_BODY, Pt(12), color=DARK_GRAY, space_before=4)
    total += 1

    # --- Slide: 例2-周期预测 ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "例题·周期结构与排布预测", "综合应用量子数+排布规则")
    y = 1.5
    add_textbox(slide, ML, y, 10, 1.0,
                "根据电子填入轨道顺序预测：\n(1) 第 8 周期共有多少种元素？\n(2) 原子核外出现第一个 6f 电子的元素原子序数是？\n(3) 第 114 号元素属于第几周期、哪一族？外围电子构型？",
                CN_BODY, Pt(14), False, DARK_GRAY)
    y += 1.4
    card = add_card(slide, ML, y, CONTENT_W, 1.5, LIGHT_GREEN)
    tf = card.text_frame; tf.word_wrap = True
    r = tf.paragraphs[0].add_run(); r.text = "  ✅ 解答要点"
    set_font(r, CN_TITLE, Pt(14), True, GREEN)
    add_paragraph(tf, "  (1) 第8周期 → 8s5g6f7d8p → 共 50 种元素（s²+g¹⁸+f¹⁴+d¹⁰+p⁶）",
                  CN_BODY, Pt(11), color=DARK_GRAY, space_before=3)
    add_paragraph(tf, "  (2) 6f 第一个电子：Z = 86(Rn) + [5g¹⁸=18] + [6f¹=1] = 105？→需按实际能级顺序",
                  CN_BODY, Pt(11), color=DARK_GRAY, space_before=2)
    add_paragraph(tf, "  (3) 114 号：第7周期、ⅣA族、[Rn]5f¹⁴6d¹⁰7s²7p²",
                  CN_BODY, Pt(11), True, DARK_BLUE, space_before=2)
    total += 1

    # --- Slide: 例3-核反应 ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "例题·核反应方程式", "2010年全国初赛试题")
    y = 1.3
    add_textbox(slide, ML, y, 10, 1.2,
                "2009年合成了第117号元素Uus。用²⁴⁹Bk轰击⁴⁸Ca靶合成，得到6个Uus原子。\n其中1个经p次α衰变得²⁷⁰Db后裂变；5个经q次α衰变得²⁸¹Rg后裂变。\n\n写出得到117号元素Uus的核反应方程式。",
                CN_BODY, Pt(13), False, DARK_GRAY)
    y += 1.8
    card = add_card(slide, ML, y, CONTENT_W, 1.5, LIGHT_GREEN)
    tf = card.text_frame; tf.word_wrap = True
    r = tf.paragraphs[0].add_run(); r.text = "  ✅ 解答"
    set_font(r, CN_TITLE, Pt(14), True, GREEN)
    add_paragraph(tf, "  ²⁴⁹Bk + ⁴⁸Ca → ²⁹⁷Uus* → ²⁹⁴Uus + 3n",
                  CN_BODY, Pt(13), True, DARK_BLUE, space_before=4)
    add_paragraph(tf, "  （²⁴⁹Bk 97号 + ⁴⁸Ca 20号 → ²⁹⁷Uus* → ²⁹⁴Uus + 3n 质量数守恒）",
                  CN_BODY, Pt(11), False, GRAY, space_before=3)
    total += 1

    # ================================================================
    # SLIDE: 总结
    # ================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid(); bg.fill.fore_color.rgb = DARK_BLUE_BG; bg.line.fill.background()
    # Title
    add_textbox(slide, ML, 0.8, 10, 0.8, "📋 本节核心速查",
                CN_TITLE, Pt(30), True, WHITE, PP_ALIGN.CENTER)
    # Summary boxes
    summaries = [
        ("四量子数", "n:层 l:形 m:向 ms:自旋"),
        ("排布三原则", "能量最低·Pauli·Hund"),
        ("充填≠失电", "填:4s先于3d, 失:4s先于3d"),
        ("两反常", "Cr/Cu半满全满·电离能ⅡA>ⅢA"),
        ("两守恒", "核反应:质量数+电荷数守恒"),
        ("一核心思维", "放弃轨道→接受概率云"),
    ]
    y = 2.2
    for i, (title, desc) in enumerate(summaries):
        x = ML + (i % 3) * 4.1
        y_pos = y + (i // 3) * 2.0
        card = add_card(slide, x, y_pos, 3.7, 1.5, LIGHT_GRAY)
        tf = card.text_frame; tf.word_wrap = True
        r = tf.paragraphs[0].add_run(); r.text = f"  {title}"
        set_font(r, CN_TITLE, Pt(15), True, DARK_BLUE)
        add_paragraph(tf, f"  {desc}", CN_BODY, Pt(12), color=DARK_GRAY, space_before=6)

    # Footer
    add_textbox(slide, ML, 6.5, 10, 0.4,
                "📖 第八讲原子结构（三）  |  📄 原子结构-超级充实版（自学完整）  |  🧠 教学洞察-原子结构",
                CN_BODY, Pt(9), False, GRAY, PP_ALIGN.CENTER)
    total += 1

    # ================================================================
    # SAVE
    # ================================================================
    # Add page numbers
    for i, slide in enumerate(prs.slides):
        if i < len(prs.slides):
            try:
                add_page_number(slide, i + 1, total)
            except:
                pass

    os.makedirs(OUTPUT_DIR, exist_ok=True)
    prs.save(OUTPUT_FILE)
    return total


if __name__ == '__main__':
    n = build_presentation()
    print("[OK] PPT generated successfully! %d slides" % n)
    print("[OK] File: " + repr(OUTPUT_FILE))
